import io
import logging
import subprocess
import tempfile
from pathlib import Path

import fitz
from pptx import Presentation as PresentationFactory
from pptx.presentation import Presentation
from pptx.util import Inches

from typst2pptx.config import ConvertConfig

logger = logging.getLogger(__name__)


def _compile_typst(typ_path: Path, config: ConvertConfig) -> bytes:
    """
    Компилирует .typ файл в PDF через typst CLI.

    Args:
        typ_path (Path): Путь к исходному .typ файлу.
        config (ConvertConfig): Конфигурация конвертации.

    Returns:
        bytes: Содержимое PDF файла.

    Raises:
        subprocess.CalledProcessError: Если typst завершился с ошибкой.
        FileNotFoundError: Если входной файл не существует.
    """
    if not typ_path.exists():
        raise FileNotFoundError(f"Input file not found: {typ_path}")

    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        tmp_path = Path(tmp.name)

    try:
        result = subprocess.run(
            [config.typst_bin, "compile", str(typ_path), str(tmp_path)],
            capture_output=True,
            text=True,
        )
        if result.returncode != 0:
            raise subprocess.CalledProcessError(
                result.returncode,
                result.args,
                output=result.stdout,
                stderr=result.stderr,
            )
        logger.debug("typst compiled: %s", typ_path)
        return tmp_path.read_bytes()
    finally:
        tmp_path.unlink(missing_ok=True)


def _pdf_to_images(
    pdf_bytes: bytes, dpi: int
) -> list[tuple[bytes, tuple[int, int]]]:
    """
    Конвертирует страницы PDF в PNG изображения.

    Args:
        pdf_bytes (bytes): Содержимое PDF файла.
        dpi (int): Разрешение рендера.

    Returns:
        list[tuple[bytes, tuple[int, int]]]: Список пар (png_bytes, (width_px, height_px)).
    """
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    scale = dpi / 72
    matrix = fitz.Matrix(scale, scale)

    result = []
    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=matrix)
        png_bytes = pix.tobytes("png")
        result.append((png_bytes, (pix.width, pix.height)))
        logger.debug("rendered page %d/%d", i + 1, len(doc))

    return result


def _build_pptx(
    pages: list[tuple[bytes, tuple[int, int]]],
    dpi: int,
) -> Presentation:
    """
    Собирает PPTX из списка PNG изображений.

    Args:
        pages (list[tuple[bytes, tuple[int, int]]]): Страницы как PNG + размеры в пикселях.
        dpi (int): Разрешение, использованное при рендере.

    Returns:
        Presentation: Готовый объект презентации.
    """
    _, (first_w, first_h) = pages[0]
    slide_w = Inches(first_w / dpi)
    slide_h = Inches(first_h / dpi)

    prs = PresentationFactory()
    prs.slide_width = slide_w
    prs.slide_height = slide_h

    blank_layout = prs.slide_layouts[6]

    for i, (png_bytes, (w_px, h_px)) in enumerate(pages):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            io.BytesIO(png_bytes),
            left=0,
            top=0,
            width=Inches(w_px / dpi),
            height=Inches(h_px / dpi),
        )
        logger.debug("added slide %d/%d", i + 1, len(pages))

    return prs


def convert(typ_path: Path, output_path: Path, config: ConvertConfig) -> None:
    """
    Конвертирует .typ файл в .pptx.

    Args:
        typ_path (Path): Путь к исходному .typ файлу.
        output_path (Path): Путь для сохранения .pptx файла.
        config (ConvertConfig): Конфигурация конвертации.

    Raises:
        FileNotFoundError: Если входной файл не существует.
        subprocess.CalledProcessError: Если typst завершился с ошибкой.
    """
    pdf_bytes = _compile_typst(typ_path, config)
    pages = _pdf_to_images(pdf_bytes, config.dpi)
    prs = _build_pptx(pages, config.dpi)
    prs.save(str(output_path))
    logger.info("saved %d slides → %s", len(pages), output_path)
