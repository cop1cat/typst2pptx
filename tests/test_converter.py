import subprocess
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from pptx import Presentation as PresentationFactory
from pptx.util import Inches

from typst2pptx.config import ConvertConfig
from typst2pptx.converter import (
    _build_pptx,
    _compile_typst,
    _pdf_to_images,
    convert,
)


@pytest.fixture
def config() -> ConvertConfig:
    return ConvertConfig(dpi=72, typst_bin="typst")


# --- _compile_typst ---


def test_compile_typst_file_not_found(
    tmp_path: Path, config: ConvertConfig
) -> None:
    missing = tmp_path / "missing.typ"
    with pytest.raises(FileNotFoundError, match=r"missing\.typ"):
        _compile_typst(missing, config)


def test_compile_typst_typst_error(
    tmp_path: Path, config: ConvertConfig
) -> None:
    typ_file = tmp_path / "slide.typ"
    typ_file.write_text("")

    mock_result = MagicMock()
    mock_result.returncode = 1
    mock_result.stdout = ""
    mock_result.stderr = "compile error"
    mock_result.args = ["typst", "compile", str(typ_file)]

    with (
        patch("typst2pptx.converter.subprocess.run", return_value=mock_result),
        pytest.raises(subprocess.CalledProcessError),
    ):
        _compile_typst(typ_file, config)


def test_compile_typst_success(
    tmp_path: Path, config: ConvertConfig, minimal_pdf_bytes: bytes
) -> None:
    typ_file = tmp_path / "slide.typ"
    typ_file.write_text("")

    def fake_run(cmd: list[str], **kwargs: object) -> MagicMock:
        Path(cmd[-1]).write_bytes(minimal_pdf_bytes)
        result = MagicMock()
        result.returncode = 0
        result.stdout = ""
        result.stderr = ""
        return result

    with patch("typst2pptx.converter.subprocess.run", side_effect=fake_run):
        pdf = _compile_typst(typ_file, config)

    assert pdf == minimal_pdf_bytes


def test_compile_typst_cleans_temp_file(
    tmp_path: Path, config: ConvertConfig, minimal_pdf_bytes: bytes
) -> None:
    typ_file = tmp_path / "slide.typ"
    typ_file.write_text("")
    captured_tmp: list[str] = []

    def fake_run(cmd: list[str], **kwargs: object) -> MagicMock:
        captured_tmp.append(cmd[-1])
        Path(cmd[-1]).write_bytes(minimal_pdf_bytes)
        result = MagicMock()
        result.returncode = 0
        return result

    with patch("typst2pptx.converter.subprocess.run", side_effect=fake_run):
        _compile_typst(typ_file, config)

    assert not Path(captured_tmp[0]).exists()


# --- _pdf_to_images ---


def test_pdf_to_images_single_page(minimal_pdf_bytes: bytes) -> None:
    pages = _pdf_to_images(minimal_pdf_bytes, dpi=72)
    assert len(pages) == 1
    png_bytes, (w, h) = pages[0]
    assert isinstance(png_bytes, bytes)
    assert png_bytes[:8] == b"\x89PNG\r\n\x1a\n"  # PNG magic bytes
    assert w > 0 and h > 0


def test_pdf_to_images_dpi_scaling(minimal_pdf_bytes: bytes) -> None:
    pages_72 = _pdf_to_images(minimal_pdf_bytes, dpi=72)
    pages_144 = _pdf_to_images(minimal_pdf_bytes, dpi=144)
    _, (w72, h72) = pages_72[0]
    _, (w144, h144) = pages_144[0]
    assert w144 == pytest.approx(w72 * 2, rel=0.05)
    assert h144 == pytest.approx(h72 * 2, rel=0.05)


def test_pdf_to_images_multi_page() -> None:
    import fitz

    doc = fitz.open()
    for _ in range(3):
        doc.new_page(width=595, height=842)
    pdf_bytes = doc.tobytes()

    pages = _pdf_to_images(pdf_bytes, dpi=72)
    assert len(pages) == 3


# --- _build_pptx ---


def test_build_pptx_single_slide(minimal_pdf_bytes: bytes) -> None:
    pages = _pdf_to_images(minimal_pdf_bytes, dpi=72)
    prs = _build_pptx(pages, dpi=72)
    assert len(prs.slides) == 1


def test_build_pptx_slide_dimensions(minimal_pdf_bytes: bytes) -> None:
    dpi = 72
    pages = _pdf_to_images(minimal_pdf_bytes, dpi=dpi)
    _, (w_px, h_px) = pages[0]
    prs = _build_pptx(pages, dpi=dpi)
    assert prs.slide_width == Inches(w_px / dpi)
    assert prs.slide_height == Inches(h_px / dpi)


def test_build_pptx_multiple_slides() -> None:
    import fitz

    doc = fitz.open()
    for _ in range(3):
        doc.new_page(width=595, height=842)
    pdf_bytes = doc.tobytes()

    pages = _pdf_to_images(pdf_bytes, dpi=72)
    prs = _build_pptx(pages, dpi=72)
    assert len(prs.slides) == 3


# --- convert ---


def test_convert_creates_pptx(tmp_path: Path, minimal_pdf_bytes: bytes) -> None:
    typ_file = tmp_path / "slide.typ"
    typ_file.write_text("")
    output = tmp_path / "out.pptx"
    config = ConvertConfig(dpi=72)

    with patch(
        "typst2pptx.converter._compile_typst", return_value=minimal_pdf_bytes
    ):
        convert(typ_file, output, config)

    assert output.exists()
    prs = PresentationFactory(str(output))
    assert len(prs.slides) == 1
